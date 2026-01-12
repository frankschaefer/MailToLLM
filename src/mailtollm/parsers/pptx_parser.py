from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from mailtollm.models.schema import AttachmentContent


def parse_pptx(path: Path, attachment_id: str) -> AttachmentContent:
    prs = Presentation(path)
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_parts.append(text)
    return AttachmentContent(id=attachment_id, text="\n".join(text_parts).strip())
